import open3d as o3d
import numpy as np
import h5py
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
import time

h5_file = 'Rev9 Steering Column Laminate.h5'

f = h5py.File(h5_file, 'r')
meshes = [item for item in f['composite_cae/meshes']]
plies = [item for item in f['composite_cae/components']]

morient = []
thickness = []
materials = []
orientations = []
names = []
layers = []


for ply in plies:
    morient.append(round(f['composite_cae/components/'+ply+'/data_map/matorient_0_angles'][0]))
    thickness.append(f['composite_cae/components/'+ply+'/data_map/thicknesses'][0])
    materials.append(f['composite_cae/components/'+ply+'/material'].attrs['TITLE'].decode('UTF-8'))
    orientations.append(f['composite_cae/components/'+ply].attrs['orientation'][0])
    names.append(f['composite_cae/components/'+ply].attrs['TITLE'].decode('UTF-8'))
    layers.append(f['composite_cae/components/'+ply].attrs['layer_id'][0])

s_layers = layers.copy()
s_layers.sort()

def capture_image(vis,mesh):
        image = vis.capture_screen_float_buffer(do_render=False)
        #plt.rcParams['figure.figsize'] = [20, 20]
        plt.figure(figsize=(18, 18))
        plt.imshow(np.asarray(image))
        plt.axis('off')
        #plt.figure(figsize=(20, 20))
        #plt.figure(num=1,dpi=(1000))
        plt.savefig(mesh,dpi=500)
        
        return False
    
def custom_draw_geometry_with_custom_fov(pcd, fov_step,mesh):
    vis = o3d.visualization.Visualizer()
    vis.create_window(window_name=str(mesh),visible=True)
    vis.add_geometry(pcd)
    if len(allMs) > 0:
            [vis.add_geometry(o3d.geometry.LineSet.create_from_triangle_mesh(m)) for m in allMs]
    #vis.add_geometry(pcd2)
    ctr = vis.get_view_control()
    ctr.rotate(10.0, 0.0)
    ctr.scale(2.0)
    print("Field of view (before changing) %.2f" % ctr.get_field_of_view())
    ctr.change_field_of_view(step=fov_step)
    print("Field of view (after changing) %.2f" % ctr.get_field_of_view())
    #o3d.visualization.RenderOption().light_on = True
    vis.get_render_option().light_on = True
    vis.get_render_option().line_width = 0.0000000001
    vis.get_render_option().show_coordinate_frame = False
    vis.get_render_option().mesh_show_back_face = True


    
    
    vis.run()
    #vis.capture_screen_image('true.png', do_render=True)
    #vis.capture_screen_image('false.png', do_render=False)
    #time.sleep(1)
    capture_image(vis,mesh)
##    image = vis.capture_screen_float_buffer()
##    plt.figure(figsize=(20, 20),dpi=(500))
##    plt.imshow(np.asarray(image))
##    plt.axis('off')
##    #plt.figure(figsize=(20, 20))
##    #plt.figure(num=1,dpi=(1000))
##    plt.show()
##    plt.savefig(mesh,dpi=500)
    vis.close()
    vis.destroy_window()

allMs = []

print("Creating mesh renders")

for i,mesh in enumerate(meshes):
    
    print(str(100*(i/len(meshes))) + "% complete")
    
    nodes = f['composite_cae/meshes/'+mesh+'/nodes']
    
    elements = f['composite_cae/meshes/'+str(mesh)+'/element_nodes']

##    x = np.linspace(-3, 3, 401)
##    mesh_x, mesh_y = np.meshgrid(x, x)
##    z = np.sinc((np.power(mesh_x, 2) + np.power(mesh_y, 2)))
##    z_norm = (z - z.min()) / (z.max() - z.min())
##    z_norm = 10
##    xyz = np.zeros((np.size(mesh_x), 3))
##    xyz[:, 0] = np.reshape(mesh_x, -1)
##    xyz[:, 1] = np.reshape(mesh_y, -1)
##    xyz[:, 2] = np.reshape(z, -1)
##    print('xyz')
##    print(xyz)

    pcd = o3d.geometry.PointCloud()
    pcd.points = o3d.utility.Vector3dVector(nodes)

    M = o3d.geometry.TriangleMesh()
    M.vertices = o3d.utility.Vector3dVector(nodes)
    M.triangles = o3d.utility.Vector3iVector(elements)
    M.paint_uniform_color([1, 0.706, 0])
    M.compute_vertex_normals()
    allMs.append(M)
    
            
    #M.get_non_manifold_edges(allow_boundary_edges=True)

    
    #custom_draw_geometry_with_custom_fov(pcd, -90.0)
[custom_draw_geometry_with_custom_fov(m, -90.0,meshes[i]) for i,m in enumerate(allMs)]
    #o3d.visualization.draw_geometries([M],mesh_show_back_face=True)
##    alpha = 0.03
##    mesh = o3d.geometry.TriangleMesh.create_from_point_cloud_alpha_shape(pcd, alpha)
##    mesh.compute_vertex_normals()
##    o3d.visualization.draw_geometries([mesh], mesh_show_back_face=True)
##    
##    o3d.visualization.draw_geometries([pcd])



prs = Presentation('Template.pptx')

slide_master = prs.slide_master
master_shapes = slide_master.shapes
master_shape = slide_master.shapes[1]
#master_shape.text = "Set name of the presentation"
m_font = master_shape.text_frame.paragraphs[0].font
m_font.name = 'Josefin Sans'
m_font.size = Pt(9)
m_font.color.rgb = RGBColor(0x4,  0x1e, 0x42)

pres_name = ' '.join(str("Imported Laminate from " + h5_file).upper())
master_shape.text_frame.paragraphs[0].text = ' '.join(str("Imported Laminate from " + h5_file).upper())

title_slide_layout = prs.slide_layouts[0]
title_slide_layout.shapes[4].text = pres_name
t_font = title_slide_layout.shapes[4].text_frame.paragraphs[0].font
t_font.name = 'Josefin Sans'
t_font.size = Pt(9)
t_font.color.rgb = RGBColor(0xff,  0xff, 0xff)

#slide = prs.slides.add_slide(title_slide_layout)
slide = prs.slides[0]
title = slide.shapes.title
#subtitle = slide.placeholders[1]
third_shape = slide.shapes[2]


title.text = "The Title"
#subtitle.text = "Imported Laminate from " + h5_file
slide.shapes[2].text = "Imported Laminate from " + h5_file
print("100% complete")
print("Rotate view then close window to save image for each ply.")


def dropSlides(slidesToKeep, prs):
    indexesToRemove = [x for x in range(1, len(prs.slides._sldIdLst)+1) if x not in slidesToKeep]

    for i, slide in enumerate(prs.slides):
        id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}

        if i+1 in indexesToRemove:
            slide_id = slide.slide_id

            prs.part.drop_rel(id_dict[slide_id][1])
            del prs.slides._sldIdLst[id_dict[slide_id][0]]

    return prs
            
prs = dropSlides([1], prs)

## start of ply table
blank_slide_layout = prs.slide_layouts[8]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes

textbox = slide.shapes[0]
sp = textbox.element
sp.getparent().remove(sp)

title = slide.shapes.title
sp = title.element
sp.getparent().remove(sp)



rows = len(meshes) +1
cols = 5
left = Inches(0.5)
top = Inches(0.5)
width = Inches(6.0)
height = Inches(0.01)

table = shapes.add_table(rows, cols, left, top, width, height).table
# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(1.0)
table.columns[2].width = Inches(2.0)
table.columns[3].width = Inches(2.0)
table.columns[4].width = Inches(2.0)

# write column headings
table.cell(0, 0).text = 'Title'
table.cell(0, 1).text = 'Layer'
table.cell(0, 2).text = 'Material'
table.cell(0, 3).text = 'Orientation'
table.cell(0, 4).text = 'Thickness'

# write body cells
for n,lay in enumerate(s_layers):
    i = layers.index(lay)
    table.cell(n+1, 0).text = names[i]
    table.cell(n+1, 1).text = str(layers[i])
    table.cell(n+1, 2).text = materials[i]
    table.cell(n+1, 3).text = str(orientations[i]) + " degree"
    table.cell(n+1, 4).text = str(thickness[i])+" mm"

def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell

for cell in iter_cells(table):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(8)
            
## end of ply table 

pixels_height = 2400
pixels_width = 3200

for lay in s_layers:
        i = layers.index(lay)
#for i in range(0,len(meshes)):
        
        img_path = meshes[i]+'.png'
        blank_slide_layout = prs.slide_layouts[8]
        slide = prs.slides.add_slide(blank_slide_layout)
        textbox = slide.shapes[0]
        sp = textbox.element
        sp.getparent().remove(sp)

        title = slide.shapes.title
        sp = title.element
        sp.getparent().remove(sp)

        
        left = top = Inches(0.5)
        pic = slide.shapes.add_picture(img_path, left, top,width=Inches(6*pixels_width/pixels_height),height=Inches(6))
        pic.crop_top = 0.32
        pic.crop_bottom  = 0.32
        pic.crop_left  = 0.32
        pic.crop_right  = 0.32

        #left = Inches(5)
        #height = Inches(5.5)
        #pic = slide.shapes.add_picture(img_path, left, top, height=height)
        shapes = slide.shapes
        #shapes.title.text = 'Adding a Table'

        rows = 2
        cols = 5
        left = Inches(0.5)
        top = Inches(6.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(1.0)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(2.0)
        table.columns[4].width = Inches(2.0)

        # write column headings
        table.cell(0, 0).text = 'Title'
        table.cell(0, 1).text = 'Layer'
        table.cell(0, 2).text = 'Material'
        table.cell(0, 3).text = 'Orientation'
        table.cell(0, 4).text = 'Thickness'

        # write body cells
        table.cell(1, 0).text = names[i]
        table.cell(1, 1).text = str(layers[i])
        table.cell(1, 2).text = materials[i]
        table.cell(1, 3).text = str(orientations[i]) + " degree"
        table.cell(1, 4).text = str(thickness[i])+" mm"

prs.save('PlyBook_210721.pptx')
